import pytest
from pathlib import Path

pytest.importorskip("pptx")
from pptx import Presentation

from bot.utils.slide_parser_clean import parse_file


def test_parse_pptx_creates_and_reads(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Тестовый заголовок"
    prs.save(str(pptx_path))

    text = parse_file(str(pptx_path))
    assert "Тестовый заголовок" in text


def test_parse_unsupported_extension(tmp_path):
    f = tmp_path / "file.txt"
    f.write_text("hello")
    text = parse_file(str(f))
    assert "hello" in text
import os
from pathlib import Path
from pptx import Presentation
import os
import pytest
from pathlib import Path

pytest.importorskip("pptx")
from pptx import Presentation

from bot.utils.slide_parser_clean import parse_file


def test_parse_pptx_creates_and_reads(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Тестовый заголовок"
    prs.save(str(pptx_path))

    text = parse_file(str(pptx_path))
    assert "Тестовый заголовок" in text


def test_parse_unsupported_extension(tmp_path):
    f = tmp_path / "file.txt"
    f.write_text("hello")
    text = parse_file(str(f))
    assert "hello" in text
