import os
from pptx import Presentation
import pdfplumber


def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shp in slide.shapes:
            if hasattr(shp, 'text'):
                t = shp.text.strip()
                if t:
                    parts.append(t)
        slide_text = "\n".join(parts)
        texts.append(f"=== СЛАЙД {i} ===\n{slide_text}")
    return "\n".join(texts)


def parse_pdf(path: str) -> str:
    texts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            texts.append(f"=== СЛАЙД {i} ===\n{text}")
    return "\n".join(texts)


def parse_file(path: str) -> str:
    if path.lower().endswith('.pptx'):
        return parse_pptx(path)
    if path.lower().endswith('.pdf'):
        return parse_pdf(path)
    raise ValueError('Unsupported file type')
